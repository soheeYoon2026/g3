"""One page for BAIC: what happened to CAS-A between their STEP and a CFD surface.

Deliberately plain - a working figure, not a slide. Four renders in a row with
one-line captions, a small table of measured numbers, a few sentences on what was
assumed and what was not attempted. Built locally with PIL; a PPTX wrapper with
the same image is written when python-pptx is available.
"""

import argparse
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

ap = argparse.ArgumentParser(description=__doc__)
ap.add_argument("--out", type=Path, required=True, help="output PNG")
ap.add_argument("--raw", type=Path, required=True, help="render of the raw model with holes")
ap.add_argument("--healed", type=Path, required=True, help="render of the healed STEP")
ap.add_argument("--wrapped", type=Path, required=True, help="render of the watertight wrap")
ap.add_argument("--date", default="2026-09-04")
args = ap.parse_args()


def font(size, bold=False, mono=False):
    candidates = (
        ["/usr/share/fonts/dejavu-sans-mono-fonts/DejaVuSansMono.ttf",
         "/usr/share/fonts/dejavu/DejaVuSansMono.ttf"] if mono else
        ["/usr/share/fonts/dejavu-sans-fonts/DejaVuSans-Bold.ttf"] if bold else
        ["/usr/share/fonts/dejavu-sans-fonts/DejaVuSans.ttf"]
    ) + ["/usr/share/fonts/google-noto-cjk/NotoSansCJK-Regular.ttc"]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def tile(sheet, row, col):
    x = 10 + col * 770
    y = 44 + row * 530
    return sheet.crop((x + 1, y + 34, x + 759, y + 519))


W, H = 1920, 1080
page = Image.new("RGB", (W, H), (255, 255, 255))
d = ImageDraw.Draw(page)
INK = (20, 20, 20)
GREY = (110, 110, 110)
RULE = (200, 200, 200)

d.text((70, 48), f"CAS-A cleanup trial  ({args.date})", fill=INK, font=font(30, bold=True))
d.text((70, 92), "AOX AI Modeler + CAD kernel + wrapper. Our own choices where the intent was unknown; "
                 "treat as a reference, not a result.", fill=GREY, font=font(18))

raw = Image.open(args.raw).convert("RGB")
healed = Image.open(args.healed).convert("RGB")
wrapped = Image.open(args.wrapped).convert("RGB")

x0, y0, pw, ph, gap = 70, 140, 425, 272, 26
panels = [
    (tile(raw, 0, 0), "received — 2,847 surfaces, none sewn; coloured lines are the 46 open loops"),
    (tile(healed, 0, 0), "after automatic sewing and capping — 40 of 46 loops closed, written back to STEP"),
    (tile(wrapped, 1, 0), "underside — flat floor put at the sill line, body mirrored about y=0"),
    (tile(wrapped, 0, 0), "wrapped at 15 mm — closed surface for the solver"),
]
for i, (img, caption) in enumerate(panels):
    x = x0 + i * (pw + gap)
    page.paste(img.resize((pw, ph), Image.LANCZOS), (x, y0))
    d.rectangle([x, y0, x + pw - 1, y0 + ph - 1], outline=RULE)
    # caption wraps by hand at ~48 chars
    words, lines, cur = caption.split(), [], ""
    for w in words:
        if len(cur) + len(w) + 1 > 52:
            lines.append(cur)
            cur = w
        else:
            cur = (cur + " " + w).strip()
    lines.append(cur)
    for k, line in enumerate(lines):
        d.text((x, y0 + ph + 10 + k * 22), line, fill=INK, font=font(16))

# numbers, as a plain monospace table
ty = 500
d.line([(70, ty), (W - 70, ty)], fill=RULE, width=1)
rows = [
    ("", "received", "after sewing", "after capping", "wrapped"),
    ("surfaces / shells", "2,847 / 2,847", "2,847 / 19", "2,895 / 19", "—"),
    ("free edges", "12,724", "724", "—", "0"),
    ("open loops", "—", "46", "6 + 4 unjoined caps", "0"),
    ("triangles", "—", "—", "161k (tessellated)", "550,656"),
    ("frontal area, full car", "—", "—", "2.522 m²", "2.526 m²  (STAR-CCM+ setup: 2.52)"),
    ("volume", "—", "—", "—", "7.34 m³  (49 % of bounding box)"),
]
mono = font(17, mono=True)
cols = [70, 360, 560, 780, 1060]
for r, row in enumerate(rows):
    yy = ty + 16 + r * 27
    for c, cell in enumerate(row):
        d.text((cols[c], yy), cell, fill=GREY if r == 0 else INK, font=mono)

# notes
ny = ty + 16 + len(rows) * 27 + 22
d.line([(70, ny), (W - 70, ny)], fill=RULE, width=1)
notes = [
    "What was done automatically: sewing in three tolerance stages (1, 5, 10.5 mm), detection of the open loops, "
    "capping of the loops below 900 mm on their own edges, a check on every cap (size, position, nothing behind it), "
    "STEP round trip.",
    "What we decided ourselves, because we do not know the intent: the underbody is a flat plane at the sill line "
    "(z = 150 mm, 460 mm above the tyre contact); the model is mirrored about y = 0; the wheel rims are closed. "
    "A flat floor at that height changes the absolute Cd, so the numbers are for trends, not for comparison with a "
    "run that had the real underbody.",
    "What was not attempted: the styling panels overlap each other (the glass sits about 4 mm above the body, "
    "8,830 self-intersections in the file as delivered). The wrap covers them; a proper resolution needs the "
    "trimmed surfaces. If the cleanup has to reproduce an ANSA workflow this takes considerably longer than this trial.",
]
yy = ny + 18
for note in notes:
    words, cur = note.split(), ""
    for w in words:
        if len(cur) + len(w) + 1 > 150:
            d.text((70, yy), cur, fill=INK, font=font(16))
            yy += 23
            cur = w
        else:
            cur = (cur + " " + w).strip()
    d.text((70, yy), cur, fill=INK, font=font(16))
    yy += 23 + 12

d.text((70, H - 50), "attached: CAS-A-assumed-wrapped-15mm.stl (solver), CAS-A-assumed-half-floor.stp (CAD, half body + floor face)",
       fill=GREY, font=font(15))

args.out.parent.mkdir(parents=True, exist_ok=True)
page.save(args.out)
print(f"PNG: {args.out}  ({args.out.stat().st_size / 1e6:.1f} MB)")

try:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(args.out), 0, 0, width=prs.slide_width, height=prs.slide_height)
    pptx_path = args.out.with_suffix(".pptx")
    prs.save(pptx_path)
    print(f"PPTX: {pptx_path}  ({pptx_path.stat().st_size / 1e6:.1f} MB)")
except Exception as exc:
    print(f"PPTX 생략: {type(exc).__name__}: {exc}")
